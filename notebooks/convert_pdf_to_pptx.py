# -*- coding: utf-8 -*-
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import sys

def convert_pdf_to_pptx(pdf_path, pptx_path):
    if not os.path.exists(pdf_path):
        print("Error: The PDF file does not exist.")
        return False
        
    print("Opening PDF file...")
    doc = fitz.open(pdf_path)
    
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 6 is the index for a blank slide layout in standard templates
    blank_slide_layout = prs.slide_layouts[6] 
    
    # Create a temp directory for images
    temp_dir = "temp_pdf_pages"
    os.makedirs(temp_dir, exist_ok=True)
    
    print(f"Total pages: {len(doc)}")
    for page_num in range(len(doc)):
        print(f"Processing page {page_num + 1}/{len(doc)}...")
        page = doc.load_page(page_num)
        
        # Render page to image with higher DPI (zoom = 2.0 -> 144 DPI)
        zoom = 2.0  
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        image_path = os.path.join(temp_dir, f"page_{page_num:03d}.png")
        pix.save(image_path)
        
        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add fullscreen image
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Remove temp image file to save space
        os.remove(image_path)
        
    # Clean up temp dir
    try:
        os.rmdir(temp_dir)
    except Exception:
        pass
        
    print("Saving PPTX file...")
    prs.save(pptx_path)
    print("Conversion completed successfully!")
    return True

if __name__ == "__main__":
    # Standard file name
    pdf_file = "Slide Báo Cáo - Hệ thống gợi ý phim MovieLens.pdf"
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        
    notebooks_dir = os.path.dirname(os.path.abspath(__file__))
    pdf_path = os.path.join(notebooks_dir, pdf_file)
    
    # If the default file name is not found, search for any PDF in the notebooks folder
    if not os.path.exists(pdf_path):
        pdf_files = [f for f in os.listdir(notebooks_dir) if f.endswith(".pdf")]
        if pdf_files:
            pdf_path = os.path.join(notebooks_dir, pdf_files[0])
            pdf_file = pdf_files[0]
            
    if not os.path.exists(pdf_path):
        print("Could not find PDF file in notebooks folder. Please save it there first.")
        sys.exit(1)
        
    pptx_name = os.path.splitext(pdf_file)[0] + ".pptx"
    pptx_path = os.path.join(notebooks_dir, pptx_name)
    
    convert_pdf_to_pptx(pdf_path, pptx_path)
